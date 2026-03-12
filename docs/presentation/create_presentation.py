# docs/presentation/create_presentation.py
import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 상수 ──────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x23, 0x32)
GRAY33  = RGBColor(0x33, 0x33, 0x33)
GRAY66  = RGBColor(0x66, 0x66, 0x66)
GRAY99  = RGBColor(0x99, 0x99, 0x99)
CARD_BG = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_DD = RGBColor(0xDD, 0xDD, 0xDD)

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "onboarding-system-presentation.pptx")


# ── 헬퍼: 텍스트박스 추가 ──────────────────────────────────
def add_textbox(slide, text, left, top, width, height,
                font_size=16, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = "맑은 고딕"
    if color:
        run.font.color.rgb = color
    return txBox


# ── 헬퍼: 카드(직사각형) 추가 ─────────────────────────────
def add_card(slide, left, top, width, height, bg_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.adjustments[0] = 0.02  # 살짝만 둥글게
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = bg_color or CARD_BG
    return shape


# ── 헬퍼: 좌측 액센트 바 추가 ─────────────────────────────
def add_accent_bar(slide, card_left, card_top, card_height):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        card_left, card_top, Cm(0.25), card_height
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    return bar


# ── 헬퍼: 카드 + 액센트 바 + 제목 + 설명 한번에 ────────────
def add_feature_card(slide, left, top, width, height, title, desc_lines):
    """desc_lines: list of str — 각 항목이 별도 단락으로 표시"""
    add_card(slide, left, top, width, height)
    add_accent_bar(slide, left, top, height)

    title_left = left + Cm(1.0)
    title_width = width - Cm(1.2)
    add_textbox(slide, title,
                title_left, top + Cm(0.4), title_width, Cm(0.9),
                font_size=16, bold=True, color=NAVY)

    desc_top = top + Cm(1.5)
    desc_height = height - Cm(1.7)
    txBox = slide.shapes.add_textbox(title_left, desc_top, title_width, desc_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(desc_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.name = "맑은 고딕"
        run.font.color.rgb = GRAY33


# ── 슬라이드 1: 표지 ────────────────────────────────────────
def _slide1(prs, blank):
    slide = prs.slides.add_slide(blank)

    add_textbox(slide, "신규입사자 온보딩 관리 시스템",
                Cm(1.5), Cm(6.5), Cm(30.87), Cm(2.5),
                font_size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_textbox(slide, "온보딩 과정의 디지털화 및 데이터 기반 관리",
                Cm(1.5), Cm(9.5), Cm(30.87), Cm(1.5),
                font_size=18, color=GRAY66, align=PP_ALIGN.CENTER)

    add_textbox(slide, "인사기획팀 박상혁 선임  |  2026. 03",
                Cm(1.5), Cm(16.5), Cm(30.87), Cm(1.0),
                font_size=13, color=GRAY99, align=PP_ALIGN.CENTER)


# ── 슬라이드 2: 기존 워크플로우 (AS-IS) ─────────────────────
def _slide2(prs, blank):
    slide = prs.slides.add_slide(blank)

    add_textbox(slide, "기존 워크플로우 (AS-IS)",
                Cm(1.5), Cm(0.8), Cm(30.87), Cm(2.2),
                font_size=28, bold=True, color=NAVY)

    boxes = [
        ("① Excel 명단 관리",  ["신규입사자 정보", "수동 입력·관리"],   Cm(1.5)),
        ("② 버디 활동 회신",   ["사진 촬영 → 스캔", "→ 이메일 제출"],  Cm(9.87)),
        ("③ 설문 수동 진행",   ["별도 도구 진행", "결과 수동 취합"],    Cm(18.24)),
        ("④ 폴더 분산 저장",   ["yCloud 여러 폴더에", "산발적 저장"],   Cm(26.61)),
    ]

    for title, desc_lines, left in boxes:
        add_card(slide, left, Cm(6.5), Cm(6.87), Cm(5.0))
        add_textbox(slide, title,
                    left + Cm(0.3), Cm(7.0), Cm(6.3), Cm(1.0),
                    font_size=14, bold=True, color=NAVY)
        txBox = slide.shapes.add_textbox(
            left + Cm(0.3), Cm(8.2), Cm(6.3), Cm(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(desc_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(12)
            run.font.name = "맑은 고딕"
            run.font.color.rgb = GRAY33

    # 화살표 3개
    for left in [Cm(8.37), Cm(16.74), Cm(25.11)]:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, Cm(8.6), Cm(1.5), Cm(0.5))
        arr.fill.solid()
        arr.fill.fore_color.rgb = NAVY
        arr.line.fill.background()


# ── 슬라이드 3: 문제점 ──────────────────────────────────────
def _slide3(prs, blank):
    slide = prs.slides.add_slide(blank)

    add_textbox(slide, "문제점",
                Cm(1.5), Cm(0.8), Cm(30.87), Cm(2.2),
                font_size=28, bold=True, color=NAVY)

    cards = [
        {
            "left": Cm(1.5),
            "title": "HR 모니터링 공백",
            "lines": [
                "멘토링·OJT가 현업에 전적으로 일임",
                "적응 상태 파악 채널 없음",
                "문제 조기 감지·개입 불가 → 조기 이탈 리스크",
            ],
        },
        {
            "left": Cm(16.87),
            "title": "데이터 분산으로 인한 관리 비효율",
            "lines": [
                "버디 사진·설문 결과·완료 여부 폴더별 분산",
                "전체 현황 파악 시 수동 취합 필요",
                "실시간 모니터링 불가",
            ],
        },
    ]

    for c in cards:
        add_card(slide, c["left"], Cm(3.3), Cm(14.5), Cm(11.0))
        add_accent_bar(slide, c["left"], Cm(3.3), Cm(11.0))
        add_textbox(slide, c["title"],
                    c["left"] + Cm(1.0), Cm(3.8), Cm(13.0), Cm(1.0),
                    font_size=16, bold=True, color=NAVY)
        txBox = slide.shapes.add_textbox(
            c["left"] + Cm(1.0), Cm(5.2), Cm(13.0), Cm(8.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(c["lines"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"• {line}"
            run.font.size = Pt(13)
            run.font.name = "맑은 고딕"
            run.font.color.rgb = GRAY33


# ── 슬라이드 4: 시스템 소개 ─────────────────────────────────
def _slide4(prs, blank):
    slide = prs.slides.add_slide(blank)

    add_textbox(slide, "시스템 소개",
                Cm(1.5), Cm(0.8), Cm(30.87), Cm(2.2),
                font_size=28, bold=True, color=NAVY)

    items = [
        ("이름",      "YURA 온보딩 관리 시스템"),
        ("접속 방식", "웹 기반, 별도 설치 없이 브라우저 접속"),
        ("계정 체계", "신규입사자·HR 관리자 계정 분리 운영"),
        ("대상 구분", "신입공채 (3개월·3차 설문)  /  경력공채 (1개월·1차 설문)"),
    ]

    for i, (label, value) in enumerate(items):
        top = Cm(4.5 + i * 3.0)
        add_textbox(slide, label,
                    Cm(1.5), top, Cm(7.0), Cm(1.2),
                    font_size=16, bold=True, color=NAVY)
        add_textbox(slide, value,
                    Cm(9.0), top, Cm(23.37), Cm(1.2),
                    font_size=16, color=GRAY33)
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Cm(1.5), top + Cm(1.5), Cm(30.87), Cm(0.03))
        sep.fill.solid()
        sep.fill.fore_color.rgb = GRAY_DD
        sep.line.fill.background()


# ── 슬라이드 5: 주요 기능 (신규입사자) ──────────────────────
def _slide5(prs, blank):
    slide = prs.slides.add_slide(blank)

    add_textbox(slide, "주요 기능 — 신규입사자",
                Cm(1.5), Cm(0.8), Cm(30.87), Cm(2.2),
                font_size=28, bold=True, color=NAVY)

    features = [
        ("온보딩 프로그램 디지털화",
         ["버디 활동 6가지 모바일 이미지 첨부, 진행 타임라인 시각화"]),
        ("설문조사 자동화",
         ["신입 3차·경력 1차, 기간 자동 관리, 주관식+척도 23문항"]),
        ("AI 감성분석 + 이메일 자동 생성",
         ["주관식 응답에서 속성별 감정 추출, 멘토·팀장 대상 이메일 초안 생성"]),
    ]

    for (title, lines), top in zip(features, [Cm(3.3), Cm(7.6), Cm(11.9)]):
        add_feature_card(slide,
                         left=Cm(1.5), top=top,
                         width=Cm(30.87), height=Cm(3.8),
                         title=title, desc_lines=lines)


# ── 슬라이드 6: 관리자 대시보드 ① 현황 모니터링 ─────────────
def _slide6(prs, blank):
    slide = prs.slides.add_slide(blank)

    add_textbox(slide, "관리자 대시보드  ①  현황 모니터링",
                Cm(1.5), Cm(0.8), Cm(30.87), Cm(2.2),
                font_size=28, bold=True, color=NAVY)

    features = [
        ("온보딩 현황",
         ["KPI 카드 4개: 전체 입사자·완료·이번주 설문마감·마감 임박자",
          "개인별 진행률 프로그레스 바"],
         Cm(3.3)),
        ("마감 알림",
         ["마감 임박자 클릭 → 대상자 팝업",
          "D-1 이하 빨간색, D-2~3 주황색 강조"],
         Cm(9.9)),
    ]

    for title, lines, top in features:
        add_feature_card(slide,
                         left=Cm(1.5), top=top,
                         width=Cm(30.87), height=Cm(6.0),
                         title=title, desc_lines=lines)

    add_textbox(slide, "관리자 기능 1/2",
                Cm(1.5), Cm(17.8), Cm(10.0), Cm(0.8),
                font_size=11, color=GRAY99)


# ── 슬라이드 7: 관리자 대시보드 ② 운영 관리 ────────────────
def _slide7(prs, blank):
    slide = prs.slides.add_slide(blank)

    add_textbox(slide, "관리자 대시보드  ②  운영 관리",
                Cm(1.5), Cm(0.8), Cm(30.87), Cm(2.2),
                font_size=28, bold=True, color=NAVY)

    features = [
        ("멘토/버디 관리",
         ["직원별 멘토·버디 배정 팝업",
          "신입·경력 구분 안내메일 HTML 자동 생성"],
         Cm(3.3)),
        ("직원 관리",
         ["CSV 일괄 업로드 또는 수기 입력(엑셀 붙여넣기)",
          "등록 즉시 계정 자동 생성"],
         Cm(7.6)),
        ("공지사항",
         ["인라인 편집, PDF 첨부·교체",
          "중요 공지 상단 고정"],
         Cm(11.9)),
    ]

    for title, lines, top in features:
        add_feature_card(slide,
                         left=Cm(1.5), top=top,
                         width=Cm(30.87), height=Cm(3.8),
                         title=title, desc_lines=lines)

    add_textbox(slide, "관리자 기능 2/2",
                Cm(1.5), Cm(17.8), Cm(10.0), Cm(0.8),
                font_size=11, color=GRAY99)


# ── 슬라이드 8: Q&A ─────────────────────────────────────────
def _slide8(prs, blank):
    slide = prs.slides.add_slide(blank)

    add_textbox(slide, "Q&A",
                Cm(1.5), Cm(4.0), Cm(30.87), Cm(3.0),
                font_size=48, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "온보딩의 전 과정을 하나의 시스템에서 관리하고,\nAI로 이탈 리스크를 조기에 감지합니다.",
                Cm(1.5), Cm(8.0), Cm(30.87), Cm(4.0),
                font_size=22, bold=True, color=GRAY33, align=PP_ALIGN.CENTER)

    add_textbox(slide, "감사합니다",
                Cm(1.5), Cm(15.5), Cm(30.87), Cm(1.5),
                font_size=16, color=GRAY99, align=PP_ALIGN.CENTER)


# ── 메인 ────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = Cm(33.87)
    prs.slide_height = Cm(19.05)
    blank = prs.slide_layouts[6]

    _slide1(prs, blank)
    _slide2(prs, blank)
    _slide3(prs, blank)
    _slide4(prs, blank)
    _slide5(prs, blank)
    _slide6(prs, blank)
    _slide7(prs, blank)
    _slide8(prs, blank)

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"저장 완료: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
