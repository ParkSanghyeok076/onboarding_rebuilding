# docs/presentation/test_presentation.py
import os
import pytest
from pptx import Presentation

PPTX_PATH = os.path.join(os.path.dirname(__file__), "onboarding-system-presentation.pptx")


def test_file_exists():
    assert os.path.exists(PPTX_PATH), f"파일 없음: {PPTX_PATH}"


def test_slide_count():
    prs = Presentation(PPTX_PATH)
    assert len(prs.slides) == 8, f"슬라이드 수 오류: {len(prs.slides)}"


def test_slide1_has_title_text():
    prs = Presentation(PPTX_PATH)
    slide1 = prs.slides[0]
    texts = [shape.text_frame.text for shape in slide1.shapes if shape.has_text_frame]
    combined = " ".join(texts)
    assert "신규입사자 온보딩 관리 시스템" in combined


def test_slide8_has_qa_label():
    prs = Presentation(PPTX_PATH)
    slide8 = prs.slides[7]
    texts = [shape.text_frame.text for shape in slide8.shapes if shape.has_text_frame]
    combined = " ".join(texts)
    assert "Q&A" in combined
